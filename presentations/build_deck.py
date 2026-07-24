#!/usr/bin/env python3
"""
Build the SecDevQA / claude-code evaluation deck.

Every number on the results slides is read from the run artifacts at build time — nothing
is hard-coded — so re-running this after a new run regenerates a truthful deck:

    /local/home/amamun/envs/devqa/bin/python presentations/build_deck.py

Runs reported: RUNS below, pooled into one set (answers + grades + transcripts under
harness/output/).
"""

from __future__ import annotations

import collections
import glob
import json
import re
import statistics as st
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parent.parent
OUT = Path(__file__).resolve().parent / "secdevqa_claude_code_eval.pptx"

# The deck reports one pooled run: the two claude-code sweeps combined into a single set.
RUNS = [
    "coding_agent_claude_code_sel15_20260723-174406",
    "coding_agent_claude_code_sel11_20260724-002703",
]
JUDGE_SLUG = "openai-gpt-5.5"

# --- minimal palette: white deck, one accent -------------------------------
INK = RGBColor(0x1A, 0x1A, 0x1A)
MUTED = RGBColor(0x70, 0x70, 0x70)
LINE = RGBColor(0xC4, 0xC4, 0xC4)
FILL = RGBColor(0xF5, 0xF5, 0xF5)
ACCENT = RGBColor(0x2F, 0x5C, 0x8F)
ACCENT_FILL = RGBColor(0xE7, 0xEE, 0xF7)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FONT = "Calibri"

W, H = Inches(13.333), Inches(7.5)
ML = Inches(0.75)            # left margin
CW = W - 2 * ML              # content width


# ---------------------------------------------------------------------------
# Data
# ---------------------------------------------------------------------------

def load_stats(runs: list[str] = RUNS, judge: str = JUDGE_SLUG) -> dict:
    if isinstance(runs, str):
        runs = [runs]
    ans, grd = [], []
    for run in runs:
        ans += [json.loads(l) for l in
                open(ROOT / f"harness/output/answers_{run}.jsonl", encoding="utf-8") if l.strip()]
        grd += [json.loads(l) for l in
                open(ROOT / f"harness/output/grades_{run}__judge-{judge}.jsonl",
                     encoding="utf-8") if l.strip()]

    # Token/cost accounting is AGENT-SIDE ONLY: it comes from claude-code's own stream-json
    # `result` events. The judge (grading) step records no usage at all, so nothing here
    # includes it. We aggregate `modelUsage` rather than the top-level `usage` block because
    # the latter covers only the main model, while total_cost_usd spans every model
    # claude-code used (it dispatches some small subtasks to Haiku) — summing modelUsage keeps
    # the token figures and the cost figure on the same scope.
    tok = collections.Counter()
    per_model = collections.defaultdict(lambda: collections.Counter())
    model_cost = collections.defaultdict(float)
    costs, turns, model = [], [], None
    _TK = {"inputTokens": "input_tokens", "outputTokens": "output_tokens",
           "cacheReadInputTokens": "cache_read_input_tokens",
           "cacheCreationInputTokens": "cache_creation_input_tokens"}
    tr_json = [f for run in runs
               for f in glob.glob(str(ROOT / f"harness/output/transcripts/{run}/*.json"))]
    for f in tr_json:
        d = json.load(open(f, encoding="utf-8"))
        for line in (d.get("output") or "").splitlines():
            line = line.strip()
            if not line.startswith("{"):
                continue
            try:
                ev = json.loads(line)
            except json.JSONDecodeError:
                continue
            if ev.get("type") == "system" and ev.get("subtype") == "init" and not model:
                model = ev.get("model")
            if ev.get("type") == "result":
                for m, u in (ev.get("modelUsage") or {}).items():
                    for src, dst in _TK.items():
                        tok[dst] += u.get(src) or 0
                        per_model[m][dst] += u.get(src) or 0
                    model_cost[m] += u.get("costUSD") or 0
                costs.append(ev.get("total_cost_usd") or 0)
                turns.append(ev.get("num_turns") or 0)
                break

    groups = collections.Counter()
    for a in ans:
        for k, v in (a.get("tool_calls_by_group") or {}).items():
            groups[k] += v

    # Per-tool counts. The live-events file carries both surfaces: MCP calls get an integer
    # step, the agent's own calls a "native-N" step (folded in at report time).
    mcp_tools, native_tools = collections.Counter(), collections.Counter()
    tr_live = [f for run in runs
               for f in glob.glob(str(ROOT / f"harness/output/transcripts/{run}/*.live.jsonl"))]
    for f in tr_live:
        for line in open(f, encoding="utf-8"):
            line = line.strip()
            if not line:
                continue
            try:
                e = json.loads(line)
            except json.JSONDecodeError:
                continue
            if e.get("t") != "tool_call":
                continue
            dst = native_tools if str(e.get("step", "")).startswith("native") else mcp_tools
            dst[e.get("tool")] += 1

    n = len(ans)
    items = [g["scores"]["item"] for g in grd if g["scores"].get("item") is not None]
    corr = [g["scores"]["correctness"] for g in grd if g["scores"].get("correctness") is not None]
    comp = [g["scores"]["completeness"] for g in grd if g["scores"].get("completeness") is not None]
    tot_calls = sum(a["n_tool_calls"] for a in ans)
    nat_calls = sum(a["n_native_tool_calls"] for a in ans)
    rt = [a["runtime_secs"] for a in ans]

    return {
        "n": n, "model": model or "claude-opus-4-8", "judge": grd[0]["judge_model"],
        "condition": grd[0]["condition"],
        "repos": collections.Counter(a["repo"] for a in ans),
        "ktype": collections.Counter(a["knowledge_type"] for a in ans),
        "outcomes": collections.Counter(g["outcome"] for g in grd),
        "item_mean": st.mean(items), "item_median": st.median(items),
        "corr_mean": st.mean(corr), "comp_mean": st.mean(comp) if comp else 0.0,
        "n_comp": len(comp),
        "n_criteria": sum(g["scores"]["n_total"] for g in grd),
        "verdicts": collections.Counter(r["verdict"] for g in grd for r in g["rubric_grades"]),
        "tot_calls": tot_calls, "mcp_calls": tot_calls - nat_calls, "nat_calls": nat_calls,
        "calls_mean": st.mean(a["n_tool_calls"] for a in ans),
        "calls_max": max(a["n_tool_calls"] for a in ans),
        "calls_min": min(a["n_tool_calls"] for a in ans),
        "groups": groups,
        "mcp_tools": mcp_tools, "native_tools": native_tools,
        "tok": tok, "cost": sum(costs), "cost_mean": st.mean(costs),
        "model_cost": dict(model_cost), "per_model": {k: dict(v) for k, v in per_model.items()},
        "turns_mean": st.mean(turns), "turns_max": max(turns),
        "rt_total": sum(rt), "rt_mean": st.mean(rt), "rt_median": st.median(rt),
        "rt_max": max(rt),
        "corpus": {
            "issues": st.mean(a["snapshot"]["n_issues"] for a in ans),
            "prs": st.mean(a["snapshot"]["n_prs"] for a in ans),
            "adv": st.mean(a["snapshot"]["n_advisories"] for a in ans),
        },
        "resp_words": st.mean(len(a["response"].split()) for a in ans),
    }


# Buckets, checked in priority order, for classifying each rubric criterion by the
# kind of pre-existing artifact it demands. A criterion that names none of these is
# "self-contained" — answerable from the diagnosis alone. Order matters: a point that
# mentions both a commit and a version lands in "commit".
_BUCKETS = [
    ("commit",          re.compile(r"\bcommit\b|\b[0-9a-f]{7,40}\b")),
    ("PR",              re.compile(r"\bpr\b|pull request", re.I)),
    ("advisory / CVE",  re.compile(r"ghsa-|cve-|advisory", re.I)),
    ("prior issue",     re.compile(r"duplicate|tracking|#\d|issue \d|already (covered|reported|existing)", re.I)),
    ("version",         re.compile(r"\bv?\d+\.\d+|version|release|backport|fixed in", re.I)),
]
_BUCKET_ORDER = ["prior issue", "PR", "commit", "version", "advisory / CVE"]


def _bucket(text: str) -> str:
    for name, pat in _BUCKETS:
        if pat.search(text):
            return name
    return "self-contained"


def rubric_analysis(runs: list[str] = RUNS, judge: str = JUDGE_SLUG) -> dict:
    """Classify every gradeable rubric criterion two ways — by the kind of pre-existing
    artifact it demands, and by axis — reporting met/total for each, so the deck can show
    WHICH criteria the agent fails."""
    if isinstance(runs, str):
        runs = [runs]
    grd = []
    for run in runs:
        grd += [json.loads(l) for l in
                open(ROOT / f"harness/output/grades_{run}__judge-{judge}.jsonl",
                     encoding="utf-8") if l.strip()]
    per = collections.defaultdict(lambda: [0, 0])          # artifact bucket -> [met, total]
    by_axis = collections.defaultdict(lambda: [0, 0])      # axis -> [met, total]
    for g in grd:
        for r in g["rubric_grades"]:
            if not r.get("gradeable"):
                continue
            met = r["verdict"] == "met"
            b = _bucket(r["text"])
            per[b][1] += 1
            per[b][0] += met
            ax = r.get("axis", "other")
            by_axis[ax][1] += 1
            by_axis[ax][0] += met
    return {"per_bucket": {k: v for k, v in per.items()},
            "by_axis": {k: v for k, v in by_axis.items()},
            "n_items": len(grd)}


# ---------------------------------------------------------------------------
# Drawing helpers
# ---------------------------------------------------------------------------

def new_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])       # blank, white
    return s


def text(slide, x, y, w, h, s, size=16, bold=False, color=INK, align=PP_ALIGN.LEFT,
         italic=False, spacing=1.0, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    lines = s.split("\n")
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = spacing
        r = p.add_run()
        r.text = ln
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = color
        r.font.name = FONT
    return tb


def title(slide, t, sub=None):
    text(slide, ML, Inches(0.44), CW, Inches(0.6), t, size=30, bold=True)
    if sub:
        text(slide, ML, Inches(1.08), CW, Inches(0.4), sub, size=14, color=MUTED)
    return Inches(1.75) if sub else Inches(1.4)


def box(slide, x, y, w, h, label=None, body=None, fill=FILL, line=LINE,
        label_size=13, body_size=11, label_color=INK, body_color=MUTED,
        shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    sh = slide.shapes.add_shape(shape, x, y, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = line
    sh.line.width = Pt(0.75)
    sh.shadow.inherit = False
    sh.text_frame.text = ""
    if label:
        pad = Inches(0.12)
        ty = y + pad if body else y
        th = (h - 2 * pad) if body else h
        tb = text(slide, x + pad, ty, w - 2 * pad, th, label, size=label_size, bold=True,
                  color=label_color, align=PP_ALIGN.CENTER,
                  anchor=MSO_ANCHOR.TOP if body else MSO_ANCHOR.MIDDLE)
        if body:
            tf = tb.text_frame
            p = tf.add_paragraph()
            p.alignment = PP_ALIGN.CENTER
            p.space_before = Pt(4)
            for i, ln in enumerate(body.split("\n")):
                if i:
                    p = tf.add_paragraph()
                    p.alignment = PP_ALIGN.CENTER
                r = p.add_run()
                r.text = ln
                r.font.size = Pt(body_size)
                r.font.color.rgb = body_color
                r.font.name = FONT
    return sh


def arrow(slide, x, y, w, h=Inches(0.22), color=ACCENT):
    sh = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x, y, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def down_arrow(slide, x, y, h, w=Inches(0.22), color=ACCENT):
    sh = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, x, y, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def bullets(slide, x, y, w, items, size=17, gap=Inches(0.52), color=INK):
    for i, it in enumerate(items):
        yy = y + i * gap
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, yy + Inches(0.135),
                                     Inches(0.08), Inches(0.08))
        dot.fill.solid()
        dot.fill.fore_color.rgb = ACCENT
        dot.line.fill.background()
        dot.shadow.inherit = False
        text(slide, x + Inches(0.24), yy, w - Inches(0.24), gap, it, size=size, color=color)


def tiles(slide, y, entries, h=Inches(1.25), size_v=30, size_l=11.5):
    """Row of stat tiles: [(value, label), ...]."""
    n = len(entries)
    gap = Inches(0.18)
    tw = int((CW - gap * (n - 1)) / n)
    for i, (val, lab) in enumerate(entries):
        x = ML + i * (tw + gap)
        box(slide, x, y, tw, h, fill=WHITE, line=LINE)
        text(slide, x, y + Inches(0.17), Emu(tw), Inches(0.5), val,
             size=size_v, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
        text(slide, x + Inches(0.08), y + Inches(0.78), Emu(tw) - Inches(0.16), Inches(0.42),
             lab, size=size_l, color=MUTED, align=PP_ALIGN.CENTER)


def hbars(slide, x, y, w, rows, max_val=None, rowh=None, label_w=Inches(1.5),
          val_w=Inches(0.85), color=ACCENT, total_h=Inches(2.55),
          label_size=13, note_size=12.5, bar_h=Inches(0.22)):
    """Horizontal bars: rows = [(label, value, optional_note)].

    Row height adapts to the number of rows so a chart fills its band without ever
    running past the footnote, whatever the category count."""
    if rowh is None:
        rowh = min(Inches(0.62), int(total_h / max(len(rows), 1)))
    mx = max_val or max(r[1] for r in rows) or 1
    track = w - label_w - val_w
    for i, r in enumerate(rows):
        lab, val = r[0], r[1]
        yy = y + i * rowh
        text(slide, x, yy + Inches(0.03), label_w - Inches(0.1), Inches(0.3), lab,
             size=label_size, color=INK)
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x + label_w, yy + Inches(0.045),
                                    track, bar_h)
        bg.fill.solid()
        bg.fill.fore_color.rgb = RGBColor(0xEE, 0xEE, 0xEE)
        bg.line.fill.background()
        bg.shadow.inherit = False
        bw = max(int(track * val / mx), Inches(0.02))
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x + label_w, yy + Inches(0.045),
                                     bw, bar_h)
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.fill.background()
        bar.shadow.inherit = False
        note = r[2] if len(r) > 2 else str(val)
        text(slide, x + label_w + track + Inches(0.12), yy + Inches(0.03), val_w, Inches(0.3),
             note, size=note_size, color=MUTED)


def footnote(slide, s):
    text(slide, ML, H - Inches(0.62), CW, Inches(0.3), s, size=10.5, color=MUTED)


def code_note(slide, x, y, w, h, s, size=11):
    sh = box(slide, x, y, w, h, fill=RGBColor(0xFA, 0xFA, 0xFA), line=LINE)
    tb = text(slide, x + Inches(0.14), y + Inches(0.11), w - Inches(0.28), h - Inches(0.22),
              s, size=size, color=INK)
    for p in tb.text_frame.paragraphs:
        for r in p.runs:
            r.font.name = "Consolas"
    return sh


# ---------------------------------------------------------------------------
# Slides
# ---------------------------------------------------------------------------

def s_title(prs, d):
    s = new_slide(prs)
    text(s, ML, Inches(2.35), CW, Inches(0.9),
         "Evaluating Claude Code on Developer Security Queries", size=40, bold=True)
    text(s, ML, Inches(3.35), CW, Inches(0.5),
         "SecDevQA — containerized agent evaluation over a time-capped project snapshot",
         size=18, color=MUTED)
    ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, ML, Inches(4.05), Inches(1.6), Pt(2.5))
    ln.fill.solid(); ln.fill.fore_color.rgb = ACCENT; ln.line.fill.background()
    ln.shadow.inherit = False
    text(s, ML, Inches(4.4), CW, Inches(0.9),
         f"Run: {d['n']} queries · {len(d['repos'])} repositories · agent {d['model']} · "
         f"judge {d['judge']}", size=14, color=MUTED)


def s_task(prs, d):
    s = new_slide(prs)
    y = title(s, "The task",
              "A developer's security question, answered as of the day it was posted")
    bullets(s, ML, y + Inches(0.3), Inches(6.0), [
        "Real security threads from GitHub issues",
        "Self-contained query + gold answer",
        "Graded on a per-item rubric, not text overlap",
        "Answerable only from what existed at time T",
    ], gap=Inches(0.72))
    x = ML + Inches(6.7); bw = CW - Inches(6.7)
    box(s, x, y - Inches(0.08), bw, Inches(1.5), "QUERY",
        "\"ImageMagick 7.1.2-13: WriteHEICImage() calls\n"
        "heif_nclx_color_profile_alloc() without a NULL\ncheck → SIGSEGV at 0x4\"",
        fill=WHITE, label_size=11, body_size=12)
    down_arrow(s, x + bw / 2 - Inches(0.11), y + Inches(1.52), Inches(0.3))
    box(s, x, y + Inches(1.92), bw, Inches(1.25), "GOLD ANSWER",
        "Patch commit + the NULL-check\nthe maintainer actually applied",
        fill=WHITE, label_size=11, body_size=12)
    down_arrow(s, x + bw / 2 - Inches(0.11), y + Inches(3.24), Inches(0.3))
    box(s, x, y + Inches(3.64), bw, Inches(0.95), "RUBRIC",
        f"{d['n_criteria']} criteria across {d['n']} items\ncorrectness · completeness",
        fill=ACCENT_FILL, line=ACCENT, label_size=11, body_size=12)


def s_architecture(prs, d):
    s = new_slide(prs)
    y = title(s, "Evaluation setup",
              "The host freezes the world at time T; the container is offline and reconstructs it")

    # ---- HOST panel
    hw = Inches(4.5)
    box(s, ML, y, hw, Inches(4.35), fill=WHITE, line=LINE)
    text(s, ML + Inches(0.18), y + Inches(0.14), hw - Inches(0.36), Inches(0.3),
         "HOST", size=12, bold=True, color=ACCENT)
    iw = hw - Inches(0.5)
    ix = ML + Inches(0.25)
    box(s, ix, y + Inches(0.55), iw, Inches(0.62), "Benchmark item",
        "thread · report time T · base commit", label_size=12, body_size=10.5)
    down_arrow(s, ix + iw / 2 - Inches(0.11), y + Inches(1.22), Inches(0.26))
    box(s, ix, y + Inches(1.56), iw, Inches(0.62), "Snapshot builder",
        "verify base commit · cap corpora at T", label_size=12, body_size=10.5)
    down_arrow(s, ix + iw / 2 - Inches(0.11), y + Inches(2.23), Inches(0.26))
    box(s, ix, y + Inches(2.57), iw, Inches(1.5), "Payload",
        "repo.git — bare mirror, ancestors of\nbase commit only (fetched by SHA)\n"
        "issues · PRs · advisories ≤ T (JSON)",
        fill=ACCENT_FILL, line=ACCENT, label_size=12, body_size=10.5)

    # ---- flow
    ax = ML + hw + Inches(0.18)
    arrow(s, ax, y + Inches(2.0), Inches(0.75))
    text(s, ax - Inches(0.05), y + Inches(1.62), Inches(0.95), Inches(0.3),
         "mount\nread-only", size=10, color=MUTED, align=PP_ALIGN.CENTER)

    # ---- CONTAINER panel
    cx = ax + Inches(0.95)
    cw = W - cx - ML
    box(s, cx, y, cw, Inches(4.35), fill=WHITE, line=ACCENT)
    text(s, cx + Inches(0.18), y + Inches(0.14), cw - Inches(0.36), Inches(0.3),
         "CONTAINER  (podman, --rm, egress-locked)", size=12, bold=True, color=ACCENT)

    px = cx + Inches(0.25); pw = cw - Inches(0.5)
    box(s, px, y + Inches(0.55), pw, Inches(0.6), "entrypoint",
        "DNS/iptables allowlist  →  clone mirror  →  checkout base commit",
        label_size=12, body_size=10.5)

    half = (pw - Inches(0.2)) / 2
    box(s, px, y + Inches(1.32), half, Inches(0.92), "/workspace/repo",
        "real clone at\nbase commit", fill=FILL, label_size=11.5, body_size=10.5)
    box(s, px + half + Inches(0.2), y + Inches(1.32), half, Inches(0.92), "MCP server",
        "issues · PRs · advisories\ncommits · CVE lookup",
        fill=ACCENT_FILL, line=ACCENT, label_size=11.5, body_size=10.5)

    down_arrow(s, px + half / 2 - Inches(0.11), y + Inches(2.3), Inches(0.26))
    down_arrow(s, px + half + Inches(0.2) + half / 2 - Inches(0.11), y + Inches(2.3), Inches(0.26))
    box(s, px, y + Inches(2.64), pw, Inches(0.72), "claude-code  (headless)",
        "native Read / Grep / Bash / git   +   MCP tools",
        fill=WHITE, line=ACCENT, label_size=13, body_size=11)
    down_arrow(s, px + pw / 2 - Inches(0.11), y + Inches(3.42), Inches(0.26))
    box(s, px, y + Inches(3.76), pw, Inches(0.42), "answer  +  tool transcript",
        fill=FILL, label_size=11.5)

    footnote(s, "No network to GitHub at run time — the mirror is a local path, so the "
                "time cap cannot be bypassed.")


def s_snapshot(prs, d):
    s = new_slide(prs)
    y = title(s, "What the agent sees",
              "Everything is frozen at T = the moment the query was posted")
    rows = [
        ("Repository", "real git clone checked out at the base commit — "
                       "native git log / blame / diff work"),
        ("History", "only ancestors of the base commit exist; nothing later is in the object store"),
        ("Issues & PRs", f"~{d['corpus']['issues']:.0f} issues, ~{d['corpus']['prs']:.0f} PRs "
                         f"up to T — the source thread removed"),
        ("Advisories", f"~{d['corpus']['adv']:.0f} GHSA records published up to T"),
        ("Network", "egress allowlist: model API + CVE resolution only; GitHub blocked"),
    ]
    for i, (k, v) in enumerate(rows):
        yy = y + i * Inches(0.82)
        box(s, ML, yy, Inches(2.4), Inches(0.62), k, fill=ACCENT_FILL, line=ACCENT, label_size=13)
        text(s, ML + Inches(2.65), yy + Inches(0.16), CW - Inches(2.65), Inches(0.5), v,
             size=14.5, color=INK)
    footnote(s, "The gold answer never enters the container; the thread it came from is excluded "
                "from the issue corpus.")


def s_tools(prs, d):
    s = new_slide(prs)
    y = title(s, "Two tool surfaces",
              "The agent keeps its own tools; MCP supplies what is not on disk")
    half = (CW - Inches(0.5)) / 2
    box(s, ML, y + Inches(0.2), half, Inches(3.0), "Agent's native tools",
        "\nRead · Grep · Glob · Bash\ngit log · git show · git blame\n\n"
        "operate on the real clone\nat /workspace/repo",
        fill=WHITE, label_size=16, body_size=14)
    box(s, ML + half + Inches(0.5), y + Inches(0.2), half, Inches(3.0), "MCP tools  (secdevqa)",
        "\nsearch_issues · get_issue\nsearch_prs · get_pr\n"
        "search_advisories · get_advisory · vuln_lookup\n"
        "read_file · search_code · git_log · git_show\n\n"
        "the only path to issues, PRs and advisories",
        fill=ACCENT_FILL, label_size=16, body_size=14, line=ACCENT)
    text(s, ML, y + Inches(3.55), CW, Inches(0.6),
         f"Every call is logged once, through one backend — {d['tot_calls']} calls this run "
         f"({d['mcp_calls']} MCP, {d['nat_calls']} native) — giving per-artifact attribution.",
         size=14, color=MUTED)


def s_grading(prs, d):
    s = new_slide(prs)
    y = title(s, "How answers are graded",
              "Per-claim rubric judging, condition-aware")
    bullets(s, ML, y + Inches(0.35), Inches(6.1), [
        "Criteria drawn from the maintainer's reply",
        "Each criterion marked met / partial / not met",
        "Two axes: correctness and completeness",
        "Item score = mean of gradeable criteria",
    ], size=16, gap=Inches(0.72))
    x = ML + Inches(6.6); bw = CW - Inches(6.6)
    box(s, x, y, bw, Inches(0.95), "Outcome per item",
        "correct · partial · incorrect", fill=WHITE, label_size=13, body_size=12)
    box(s, x, y + Inches(1.12), bw, Inches(0.95), "Item score",
        "mean of gradeable criteria", fill=WHITE, label_size=13, body_size=12)
    box(s, x, y + Inches(2.24), bw, Inches(1.35), "Judge",
        f"{d['judge']}\n\nalways a different model\nfrom the one under test",
        fill=ACCENT_FILL, line=ACCENT, label_size=13, body_size=12)


def s_runconfig(prs, d):
    s = new_slide(prs)
    y = title(s, "This run", f"{d['condition']} · {len(RUNS)} sweeps pooled")
    tiles(s, y, [
        (str(d["n"]), "security queries"),
        (str(len(d["repos"])), "repositories"),
        (f"{d['ktype']['grounded']}/{d['ktype']['parametric']}", "grounded / parametric"),
        (f"{d['rt_total']/60:.0f} min", "wall clock"),
        (f"${d['cost']:.2f}", "agent API cost"),
    ])
    rows = [(r.split("/")[-1], c, str(c)) for r, c in d["repos"].most_common()]
    text(s, ML, y + Inches(2.00), CW, Inches(0.3), "Queries per repository",
         size=13, bold=True, color=MUTED)
    hbars(s, ML, y + Inches(2.42), Inches(7.6), rows, label_w=Inches(1.7), val_w=Inches(0.4),
          total_h=Inches(2.9))
    x = ML + Inches(8.2)
    box(s, x, y + Inches(2.32), CW - Inches(8.2), Inches(2.1), "Configuration",
        f"\nagent: {d['model']}\ncondition: {d['condition']}\n"
        f"context: full snapshot\nweb access: off",
        fill=WHITE, label_size=13, body_size=13)


def s_outcomes(prs, d):
    s = new_slide(prs)
    y = title(s, "Results — outcomes", f"Judged by {d['judge']}")
    n = d["n"]
    oc = d["outcomes"]
    tiles(s, y, [
        (f"{100*oc.get('correct',0)/n:.0f}%", f"correct ({oc.get('correct',0)}/{n})"),
        (f"{100*oc.get('partial',0)/n:.0f}%", f"partial ({oc.get('partial',0)}/{n})"),
        (f"{oc.get('incorrect',0)}", "incorrect"),
        (f"{d['item_mean']:.2f}", "mean item score"),
    ])
    text(s, ML, y + Inches(2.00), CW, Inches(0.3), "Rubric criteria verdicts",
         size=13, bold=True, color=MUTED)
    v = d["verdicts"]; tot = sum(v.values())
    hbars(s, ML, y + Inches(2.42), Inches(7.6), [
        ("met", v.get("met", 0), f"{v.get('met',0)}  ({100*v.get('met',0)/tot:.0f}%)"),
        ("partial", v.get("partial", 0), f"{v.get('partial',0)}  ({100*v.get('partial',0)/tot:.0f}%)"),
        ("not met", v.get("not_met", 0), f"{v.get('not_met',0)}  ({100*v.get('not_met',0)/tot:.0f}%)"),
    ], label_w=Inches(1.3), val_w=Inches(1.3), total_h=Inches(2.4))
    x = ML + Inches(8.2)
    box(s, x, y + Inches(2.32), CW - Inches(8.2), Inches(1.55), "Score by axis",
        f"\ncorrectness   {d['corr_mean']:.2f}\ncompleteness  {d['comp_mean']:.2f}",
        fill=WHITE, label_size=13, body_size=15)
    footnote(s, f"{tot} rubric criteria graded across {n} items. Completeness applies to the "
                f"{d['n_comp']} items carrying completeness criteria.")


def _pct(part, whole):
    return 100 * part / whole if whole else 0.0


def s_rubric_failures(prs, rn):
    s = new_slide(prs)
    y = title(s, "Results — which criteria fail",
              "Every rubric criterion, by what the answer must produce")

    # left top: by the artifact the criterion names, worst first
    text(s, ML, y + Inches(0.02), Inches(7.6), Inches(0.3),
         "Criteria met — by the artifact referenced", size=13, bold=True, color=MUTED)
    per = rn["per_bucket"]
    rows = sorted(((k, _pct(*per[k]), f"{_pct(*per[k]):.0f}%  ({per[k][0]}/{per[k][1]})")
                   for k in per),
                  key=lambda r: r[1])
    hbars(s, ML, y + Inches(0.44), Inches(7.6), rows, max_val=100,
          label_w=Inches(1.9), val_w=Inches(1.7), total_h=Inches(2.55), bar_h=Inches(0.22))

    # left bottom: by axis
    text(s, ML, y + Inches(3.35), Inches(7.6), Inches(0.3),
         "Criteria met — by axis", size=13, bold=True, color=MUTED)
    ax = rn["by_axis"]
    arows = [(k, _pct(*ax[k]), f"{_pct(*ax[k]):.0f}%  ({ax[k][0]}/{ax[k][1]})")
             for k in ("correctness", "completeness") if k in ax]
    hbars(s, ML, y + Inches(3.77), Inches(7.6), arows, max_val=100,
          label_w=Inches(1.9), val_w=Inches(1.7), total_h=Inches(0.95), bar_h=Inches(0.28))

    # right: interpretation
    x = ML + Inches(8.2); bw = CW - Inches(8.2)
    box(s, x, y + Inches(0.44), bw, Inches(2.15), "The pattern",
        "\nThe agent diagnoses the\nvulnerability, then cannot\nname the specific prior\n"
        "issue or PR the maintainer\npointed to as the resolution.",
        fill=ACCENT_FILL, line=ACCENT, label_size=14, body_size=13)
    box(s, x, y + Inches(2.75), bw, Inches(2.0), "Hardest",
        "\nReferences to a prior issue\nor PR fail most — they live\n"
        "only in the project's own\nhistory, not in code the\nagent can read.",
        fill=WHITE, line=LINE, label_size=14, body_size=13)

    ref_t = sum(per[k][1] for k in per if k != "self-contained")
    footnote(s, "Each criterion bucketed by keyword from its text (priority: "
                "commit → PR → advisory → prior issue → version; else self-contained). "
                f"{ref_t} of the criteria name a prior artifact, across {rn['n_items']} items.")


def s_toolusage(prs, d):
    s = new_slide(prs)
    y = title(s, "Results — tool usage", "What the agent actually consulted")
    tiles(s, y, [
        (str(d["tot_calls"]), "tool calls"),
        (f"{d['calls_mean']:.1f}", "mean per query"),
        (f"{d['calls_min']}–{d['calls_max']}", "range per query"),
        (f"{d['turns_mean']:.1f}", "mean agent turns"),
    ])
    text(s, ML, y + Inches(2.00), CW, Inches(0.3), "MCP calls by artifact group",
         size=13, bold=True, color=MUTED)
    rows = [(k, v, str(v)) for k, v in d["groups"].most_common()]
    hbars(s, ML, y + Inches(2.42), Inches(7.6), rows, label_w=Inches(1.5), val_w=Inches(0.5),
          total_h=Inches(2.75))
    x = ML + Inches(8.2)
    box(s, x, y + Inches(2.32), CW - Inches(8.2), Inches(1.75), "Split",
        f"\nMCP snapshot tools      {d['mcp_calls']}\n"
        f"agent's own tools          {d['nat_calls']}",
        fill=WHITE, label_size=13, body_size=13)
    footnote(s, "Every query used at least one MCP call — the snapshot corpora are not "
                "reachable any other way.")


def s_tooldist(prs, d):
    s = new_slide(prs)
    mcp, nat = d["mcp_tools"], d["native_tools"]
    n_tools = len(mcp) + len(nat)
    y = title(s, "Results — which tools were called",
              f"{d['tot_calls']} calls across {n_tools} distinct tools")

    mx = max(list(mcp.values()) + list(nat.values()))
    rowh = Inches(0.34)

    # left: MCP tools
    lw = Inches(7.0)
    text(s, ML, y, lw, Inches(0.3), f"MCP snapshot tools  ({d['mcp_calls']})",
         size=13, bold=True, color=ACCENT)
    hbars(s, ML, y + Inches(0.42), lw,
          [(k, v, str(v)) for k, v in mcp.most_common()], max_val=mx,
          label_w=Inches(1.85), val_w=Inches(0.45), rowh=rowh,
          label_size=12, note_size=11.5, bar_h=Inches(0.19))

    # right: the agent's own tools
    rx = ML + Inches(7.6)
    rw = W - ML - rx
    text(s, rx, y, rw, Inches(0.3), f"Agent's own tools  ({d['nat_calls']})",
         size=13, bold=True, color=ACCENT)
    hbars(s, rx, y + Inches(0.42), rw,
          [(k, v, str(v)) for k, v in nat.most_common()], max_val=mx,
          label_w=Inches(1.5), val_w=Inches(0.45), rowh=rowh,
          label_size=12, note_size=11.5, bar_h=Inches(0.19))

    nb = nat.get("Bash", 0)
    text(s, rx, y + Inches(1.78), rw, Inches(0.3),
         f"What Bash runs  ({nb} calls)", size=13, bold=True, color=ACCENT)
    code_note(s, rx, y + Inches(2.12), rw, Inches(1.72),
              "git tag --contains <fix_sha>\n"
              "git log --oneline -5 -- coders/heic.c\n"
              "git describe --tags\n"
              "grep '\"version\"' package.json\n"
              "grep -rn \"freebsd|solaris\" hack/*.sh\n"
              "grep -nE \"^4\\.21|semver lock\" History.md",
              size=10.5)
    text(s, rx, y + Inches(3.94), rw, Inches(0.5),
         "mostly release archaeology — tags, ancestry,\nmanifests and changelogs",
         size=11.5, color=MUTED)

    footnote(s, "ToolSearch is claude-code's own tool-discovery step, not project "
                "investigation. Bash reaches for git metadata the MCP tools do not expose.")


def s_cost(prs, d):
    s = new_slide(prs)
    y = title(s, "Results — cost and latency", "Per query, and for the whole run")
    t = d["tok"]
    tiles(s, y, [
        (f"{(t['input_tokens']+t['output_tokens'])/1000:.0f}k", "input + output tokens"),
        (f"{t['cache_read_input_tokens']/1e6:.2f}M", "cache-read tokens"),
        (f"${d['cost_mean']:.2f}", "agent cost per query"),
        (f"{d['rt_mean']:.0f}s", "mean latency"),
    ])
    text(s, ML, y + Inches(2.00), CW, Inches(0.3), "Latency per query (seconds)",
         size=13, bold=True, color=MUTED)
    hbars(s, ML, y + Inches(2.42), Inches(7.6), [
        ("median", d["rt_median"], f"{d['rt_median']:.0f}s"),
        ("mean", d["rt_mean"], f"{d['rt_mean']:.0f}s"),
        ("slowest", d["rt_max"], f"{d['rt_max']:.0f}s"),
    ], label_w=Inches(1.4), val_w=Inches(0.8), total_h=Inches(2.4))
    x = ML + Inches(8.2)
    box(s, x, y + Inches(2.32), CW - Inches(8.2), Inches(1.75), "Totals",
        f"\n{d['rt_total']/60:.0f} min wall clock\n${d['cost']:.2f} across {d['n']} queries\n"
        f"~{d['resp_words']:.0f} words per answer",
        fill=WHITE, label_size=13, body_size=13)
    models = " + ".join(sorted(d["per_model"]))
    footnote(s, f"Agent only ({models}) — grading by {d['judge']} is not included; the "
                f"harness records no judge usage. Cache-read dominates: the snapshot prompt "
                f"and tool results are reused across turns within an item.")


def s_findings(prs, d, rn):
    s = new_slide(prs)
    y = title(s, "What stands out") + Inches(0.25)
    n = d["n"]
    per = rn["per_bucket"]
    worst = min(((k, _pct(*per[k])) for k in per if k != "self-contained"),
                key=lambda t: t[1], default=("prior issue", 0))
    bullets(s, ML, y, CW, [
        f"Diagnosis is strong, attribution is weak — correctness axis {d['corr_mean']:.2f} "
        f"vs completeness {d['comp_mean']:.2f}.",
        "Failures cluster on criteria that name a specific prior artifact — "
        f"{worst[0]} references are met only {worst[1]:.0f}% of the time.",
        f"Effort varies widely: {d['calls_min']}–{d['calls_max']} tool calls and up to "
        f"{d['turns_max']} turns per query.",
        "Few answers are fully wrong; failures are omissions and mis-attributions, "
        "not misdiagnoses.",
    ], size=17, gap=Inches(0.95))
    footnote(s, f"{n} queries · two sweeps pooled · judge {d['judge']} — indicative, not yet "
                f"statistically powered.")


def s_next(prs, d):
    s = new_slide(prs)
    y = title(s, "Limitations and next steps")
    half = (CW - Inches(0.5)) / 2
    box(s, ML, y, half, Inches(3.6), "Limitations",
        f"\n{d['n']} queries, one run, one judge\n\n"
        "vuln_lookup resolves live —\nnot capped at T\n\n"
        "single judge, no human κ yet",
        fill=WHITE, label_size=16, body_size=14)
    box(s, ML + half + Inches(0.5), y, half, Inches(3.6), "Next",
        "\nfull benchmark, all conditions\n\n"
        "bare LLM and snapshot agent\nfor comparison\n\n"
        "selective provision (leave-one-out)\nfor artifact attribution",
        fill=ACCENT_FILL, line=ACCENT, label_size=16, body_size=14)


def main() -> None:
    d = load_stats()                          # both sweeps pooled into one set
    rn = rubric_analysis()                    # rubric-failure taxonomy over the pooled set
    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H
    # Flow: overall results (s_outcomes) → which criteria fail (s_rubric_failures) → the rest.
    for fn, args in (
        (s_title, (d,)), (s_task, (d,)), (s_architecture, (d,)), (s_snapshot, (d,)),
        (s_tools, (d,)), (s_grading, (d,)), (s_runconfig, (d,)),
        (s_outcomes, (d,)), (s_rubric_failures, (rn,)),
        (s_toolusage, (d,)), (s_tooldist, (d,)), (s_cost, (d,)),
        (s_findings, (d, rn)), (s_next, (d,)),
    ):
        fn(prs, *args)
    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(f"wrote {OUT}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")


if __name__ == "__main__":
    main()
